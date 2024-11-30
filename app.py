import os
import base64
import io
import logging
import traceback
import time
import gc
from flask import Flask, request, jsonify
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from concurrent.futures import ThreadPoolExecutor
import mimetypes

def create_app():
    # Existing logging and configuration setup remains the same
    app = Flask(__name__)
    
    # Enhanced configuration for large file uploads
    app.config['MAX_CONTENT_LENGTH'] = 1.1 * 1024 * 1024 * 1024  # 1.1GB max upload size
    app.config['UPLOAD_FOLDER'] = '/tmp/uploads'
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

    # Initialize resources
    executor = ThreadPoolExecutor(max_workers=10)
    from model import get_model
    summarization_model = get_model()

    def extract_text_from_document(content, file_type):
        """
        Enhanced text extraction with comprehensive error handling for multiple file types.
        
        Args:
            content (bytes/str): Document content 
            file_type (str): Type of document (pdf, docx, xlsx, etc.)
        
        Returns:
            str: Extracted text content or empty string if extraction fails
        """
        try:
            # Convert content to bytes if it's a base64 string
            if isinstance(content, str):
                try:
                    content = base64.b64decode(content)
                except Exception:
                    # If base64 decoding fails, treat as plain text
                    return content

            # Comprehensive text extraction for multiple file types
            extractors = {
                'pdf': lambda c: _extract_pdf_text(c),
                'docx': lambda c: _extract_docx_text(c),
                'doc': lambda c: _extract_docx_text(c),
                'xlsx': lambda c: _extract_excel_text(c),
                'xls': lambda c: _extract_excel_text(c),
                'pptx': lambda c: _extract_pptx_text(c),
                'ppt': lambda c: _extract_pptx_text(c),
                'txt': lambda c: c.decode('utf-8', errors='ignore'),
                'md': lambda c: c.decode('utf-8', errors='ignore'),
                'png': lambda c: _extract_image_text(c),
                'jpg': lambda c: _extract_image_text(c),
                'jpeg': lambda c: _extract_image_text(c)
            }

            # Normalize file type to lowercase
            file_type = file_type.lower()
            
            # Select appropriate extractor or use generic text extraction
            extractor = extractors.get(file_type, lambda c: c.decode('utf-8', errors='ignore'))
            
            extracted_text = extractor(content)
            return extracted_text if extracted_text else ""

        except Exception as e:
            logging.error(f"Error extracting text from {file_type} file: {str(e)}")
            return ""
        finally:
            gc.collect()

    def _extract_pdf_text(content):
        try:
            from pdfminer.high_level import extract_text
            return extract_text(io.BytesIO(content))
        except Exception as e:
            logging.error(f"PDF text extraction error: {str(e)}")
            return ""

    def _extract_docx_text(content):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join([para.text for para in doc.paragraphs if para.text])
        except Exception as e:
            logging.error(f"DOCX text extraction error: {str(e)}")
            return ""

    def _extract_excel_text(content):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell])
                    if row_text:
                        text += row_text + "\n"
            return text
        except Exception as e:
            logging.error(f"Excel text extraction error: {str(e)}")
            return ""

    def _extract_pptx_text(content):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logging.error(f"PPTX text extraction error: {str(e)}")
            return ""

    def _extract_image_text(content):
        try:
            import pytesseract
            from PIL import Image
            image = Image.open(io.BytesIO(content))
            return pytesseract.image_to_string(image) or ""
        except Exception as e:
            logging.error(f"Image text extraction error: {str(e)}")
            return ""

    @app.route('/summarize', methods=['POST'])
    def summarize_documents():
        """
        Enhanced endpoint for document summarization with robust error handling.
        """
        start_time = time.time()
        try:
            # Extract parameters from form data or JSON
            summary_depth = float(request.form.get('summary_depth', 0.3))
            language = request.form.get('language', 'english')
            user_id = request.form.get('user_id', 'default_user')

            # Determine input method (multipart form or base64 JSON)
            if request.files:
                # Multipart form file upload
                uploaded_files = request.files
                processed_documents = []

                for filename, file_storage in uploaded_files.items():
                    content = file_storage.read()
                    file_type = filename.split('.')[-1].lower()
                    
                    extracted_text = extract_text_from_document(content, file_type)
                    
                    if extracted_text.strip():
                        processed_documents.append({
                            'name': filename,
                            'content': extracted_text,
                            'type': file_type
                        })

            elif request.is_json:
                # JSON payload with base64 encoded files
                json_data = request.json
                processed_documents = []

                for file_info in json_data.get('files', []):
                    filename = file_info.get('name', 'unknown')
                    content = file_info.get('content', '')
                    file_type = filename.split('.')[-1].lower()

                    extracted_text = extract_text_from_document(content, file_type)
                    
                    if extracted_text.strip():
                        processed_documents.append({
                            'name': filename,
                            'content': extracted_text,
                            'type': file_type
                        })
            else:
                return jsonify({
                    'status': 'error', 
                    'message': 'Invalid request format. Use multipart/form-data or application/json'
                }), 400

            # Check if any documents were processed
            if not processed_documents:
                return jsonify({
                    'status': 'error', 
                    'message': 'No valid documents found for summarization'
                }), 400

            # Import summarization modules
            from summarie import generate_summary, enrich_summary_with_visuals

            # Generate summaries
            summaries = generate_summary(
                model=summarization_model,
                documents=processed_documents,
                summary_depth=summary_depth,
                language=language
            )

            # Enrich summaries with visuals
            enriched_summaries = [
                enrich_summary_with_visuals(summary) 
                for summary in summaries
            ]

            execution_time = time.time() - start_time
            
            return jsonify({
                'status': 'success',
                'summaries': enriched_summaries,
                'execution_time': execution_time
            })

        except Exception as e:
            logging.error(f"Summarization process error: {e}")
            logging.error(traceback.format_exc())
            return jsonify({
                'status': 'error', 
                'message': str(e)
            }), 500
        finally:
            gc.collect()
    
    # Feedback and Health Check routes remain the same as in the original implementation
    @app.route('/feedback', methods=['POST'])
    def feedback():
        try:
            data = request.json
            required_fields = ['summary_id', 'user_id', 'feedback', 'summary_content']
            
            if not all(field in data for field in required_fields):
                return jsonify({'error': 'Missing required fields'}), 400

            feedback_data = {
                'summary_id': data['summary_id'],
                'user_id': data['user_id'],
                'feedback': data['feedback'],
                'summary_content': data['summary_content'],
                'timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
            }

            logging.info(f"Feedback received: {json.dumps(feedback_data, indent=2)}")

            return jsonify({
                'status': 'success',
                'message': 'Feedback received successfully',
                'feedback_id': hash(str(feedback_data))
            }), 200

        except Exception as e:
            logging.error(f"Error in feedback route: {str(e)}")
            logging.error(traceback.format_exc())
            return jsonify({'error': str(e)}), 500

    @app.route('/health', methods=['GET'])
    def health_check():
        try:
            model_status = 'healthy' if summarization_model else 'unavailable'
            
            return jsonify({
                'status': 'healthy',
                'service': 'document-summarizer',
                'version': os.environ.get('APP_VERSION', '1.0.0'),
                'model_status': model_status,
                'timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
            }), 200
        except Exception as e:
            logging.error(f"Error in health check: {str(e)}")
            return jsonify({
                'status': 'unhealthy',
                'error': str(e)
            }), 500

    return app

# Create the application instance
app = create_app()

if __name__ == '__main__':
    # Load environment variables
    load_dotenv()
    
    # Run the application
    app.run(
        debug=False,
        host='0.0.0.0',
        port=int(os.environ.get('PORT', 5000)),
        threaded=True
    )
