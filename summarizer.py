import os
import requests
from datetime import datetime, timedelta
from hashlib import md5
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.probability import FreqDist
import string

# Cache and Rate Limiting Configurations
image_cache = {}
rate_limits = {
    'unsplash': {'limit': 50, 'interval': timedelta(hours=1), 'count': 0, 'reset': datetime.now()},
    'pexels': {'limit': 200, 'interval': timedelta(hours=1), 'count': 0, 'reset': datetime.now()},
    'pixabay': {'limit': 500, 'interval': timedelta(hours=1), 'count': 0, 'reset': datetime.now()}
}

def check_rate_limit(api_name):
    """ Check if API usage is within the allowed rate limits """
    if datetime.now() >= rate_limits[api_name]['reset']:
        rate_limits[api_name]['count'] = 0
        rate_limits[api_name]['reset'] = datetime.now() + rate_limits[api_name]['interval']
    
    if rate_limits[api_name]['count'] < rate_limits[api_name]['limit']:
        rate_limits[api_name]['count'] += 1
        return True
    return False

def extract_text_from_document(file):
    file_extension = file.filename.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(file)
    elif file_extension in ['doc', 'docx']:
        return extract_text_from_word(file)
    elif file_extension in ['ppt', 'pptx']:
        return extract_text_from_powerpoint(file)
    elif file_extension in ['txt', 'md']:
        return file.read().decode('utf-8')
    elif file_extension in ['jpg', 'jpeg', 'png', 'gif']:
        return extract_text_from_image(file)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_word(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_powerpoint(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_image(file):
    image = Image.open(file)
    # Use pytesseract for OCR to extract text from images
    text = pytesseract.image_to_string(image)
    return text

def summarize_document(file, summary_id, user_id, summaries_collection, model):
    content = extract_text_from_document(file)
    
    summary = model.generate_summary(content)
    
    enriched_summary = enrich_summary_with_visuals(summary)
    
    summary_data = {
        "_id": summary_id,
        "user_id": user_id,
        "document_content": content,
        "summary": enriched_summary,
        "feedback": []
    }
    summaries_collection.insert_one(summary_data)
    
    return enriched_summary

def enrich_summary_with_visuals(summary):
    keywords = extract_keywords(summary)
    images = fetch_images_from_multiple_sources(keywords)
    enriched_summary = insert_visuals_into_summary(summary, images)
    
    return enriched_summary

def extract_keywords(text):
    # Tokenize the text and remove stopwords
    tokens = word_tokenize(text.lower())
    tokens = [word for word in tokens if word.isalnum()]
    stop_words = set(stopwords.words('english') + list(string.punctuation))
    filtered_tokens = [word for word in tokens if word not in stop_words]

    # Calculate frequency distribution of the tokens
    fdist = FreqDist(filtered_tokens)
    
    # Extract the top 5 most common words as keywords
    keywords = [word for word, freq in fdist.most_common(5)]
    
    return keywords

def fetch_images_from_multiple_sources(keywords):
    """ Fetch images from Unsplash, Pexels, and Pixabay """
    images = []
    for keyword in keywords:
        if check_rate_limit('unsplash'):
            images.extend(fetch_images_from_cache_or_api(keyword, 'unsplash', fetch_image_unsplash))
        if check_rate_limit('pexels'):
            images.extend(fetch_images_from_cache_or_api(keyword, 'pexels', fetch_image_pexels))
        if check_rate_limit('pixabay'):
            images.extend(fetch_images_from_cache_or_api(keyword, 'pixabay', fetch_image_pixabay))
    return images[:3]  # Limit to 3 images to avoid cluttering the summary

def fetch_images_from_cache_or_api(keyword, api_name, fetch_func):
    cache_key = f"{api_name}_{md5(keyword.encode('utf-8')).hexdigest()}"
    if cache_key in image_cache:
        return image_cache[cache_key]
    images = fetch_func(keyword)
    image_cache[cache_key] = images
    return images

def fetch_image_unsplash(keyword):
    """ Fetch images from Unsplash """
    response = requests.get(f'https://api.unsplash.com/photos/random?query={keyword}&client_id={os.getenv("UNSPLASH_ACCESS_KEY")}')
    if response.status_code == 200:
        return [response.json()['urls']['small']]
    return []

def fetch_image_pexels(keyword):
    """ Fetch images from Pexels """
    headers = {
        'Authorization': os.getenv("PEXELS_API_KEY")
    }
    response = requests.get(f'https://api.pexels.com/v1/search?query={keyword}&per_page=1', headers=headers)
    if response.status_code == 200:
        data = response.json()
        if data['photos']:
            return [data['photos'][0]['src']['small']]
    return []

def fetch_image_pixabay(keyword):
    """ Fetch images from Pixabay """
    response = requests.get(f'https://pixabay.com/api/?key={os.getenv("PIXABAY_API_KEY")}&q={keyword}&image_type=photo&per_page=1')
    if response.status_code == 200:
        data = response.json()
        if data['hits']:
            return [data['hits'][0]['webformatURL']]
    return []

def insert_visuals_into_summary(summary, images):
    # Insert images at appropriate positions in the summary
    paragraphs = summary.split('\n\n')
    enriched_paragraphs = []
    for i, paragraph in enumerate(paragraphs):
        enriched_paragraphs.append(paragraph)
        if i < len(images):
            enriched_paragraphs.append(f'\n![Image]({images[i]})\n')
    return '\n\n'.join(enriched_paragraphs)

def feedback_improve_model(summary_id, user_id, feedback_text, summaries_collection, model):
    # Retrieve the original summary and document content
    summary_data = summaries_collection.find_one({"_id": summary_id})
    if summary_data:
        original_summary = summary_data['summary']
        original_content = summary_data['document_content']
        
        # Use the feedback to improve the model
        model.improve_model(original_content, original_summary, feedback_text)